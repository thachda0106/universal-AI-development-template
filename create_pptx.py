from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os

# ── Color Palette ───────────────────────────────────────
DARK_BG       = RGBColor(0x0D, 0x14, 0x26)   # deeper navy
CARD_BG       = RGBColor(0x15, 0x1F, 0x38)   # card surface
CARD_HOVER    = RGBColor(0x1A, 0x27, 0x45)   # slightly lighter
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # primary blue
ACCENT_GREEN  = RGBColor(0x22, 0xC5, 0x5E)   # success green
ACCENT_ORANGE = RGBColor(0xF9, 0x73, 0x16)   # warning orange
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)   # danger red
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)   # purple
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # cyan
ACCENT_YELLOW = RGBColor(0xE5, 0xB7, 0x0B)   # yellow/gold
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xE2, 0xE8, 0xF0)   # text on dark
MEDIUM_GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
BORDER_SUBTLE = RGBColor(0x33, 0x40, 0x55)   # card border
ERROR_BG      = RGBColor(0x3B, 0x0F, 0x0F)   # dark red bg
SUCCESS_BG    = RGBColor(0x0F, 0x34, 0x1A)   # dark green bg

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Reusable utilities ─────────────────────────────────

def _set_shape_shadow(shape):
    """Add a subtle drop shadow to a shape via XML"""
    # Find or create spPr element
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = shape._element.find(qn('a:spPr'))
    if spPr is None:
        spPr = etree.SubElement(shape._element, qn('p:spPr'))
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))
    outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '40000')
    outerShdw.set('dist', '25000')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'bl')
    outerShdw.set('rotWithShape', '0')
    srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '000000')
    alpha = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', '45000')

def _set_shape_border(shape, color=BORDER_SUBTLE, width=Emu(19050)):
    """Set a 1.5pt border line on a shape"""
    shape.line.color.rgb = color
    shape.line.width = width

def add_dark_bg(slide, color=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_title_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_BLUE; s.line.fill.background()

def add_slide_number(slide, n):
    box = slide.shapes.add_textbox(Inches(12.2), Inches(7.15), Inches(1), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(n); p.font.size = Pt(9); p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.RIGHT; p.font.name = 'Calibri Light'

def card(slide, l, t, w, h, color=CARD_BG, shadow=True, border=True):
    """Create a styled card shape"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if shadow: _set_shape_shadow(s)
    if border: _set_shape_border(s)
    try: s.adjustments[0] = 0.025
    except: pass
    return s

def pill(slide, l, t, w, h, color, text, size=10, text_color=WHITE, bold=True):
    """Create a small label pill"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    try: s.adjustments[0] = 0.3
    except: pass
    tf = s.text_frame; tf.word_wrap = False
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = text_color
    p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
    return s

def rect_bg(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
    """Single-line textbox with proper margins"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return box

def mtb(slide, l, t, w, h, lines, size=13, color=LIGHT_GRAY, bold_first=False, ls=1.35):
    """Multi-line textbox with paragraph spacing"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size); p.font.color.rgb = color; p.font.name = 'Calibri'
        p.space_before = Pt(size * 0.15)
        p.space_after  = Pt(size * (ls - 1) * 0.5)
        if bold_first and i == 0: p.font.bold = True
    return box

def arrow_r(slide, l, t, w=Inches(0.2), h=Inches(0.14)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = MEDIUM_GRAY; a.line.fill.background()

def section_header(slide, num, title, subtitle=None):
    add_dark_bg(slide); add_title_bar(slide); add_slide_number(slide, num)
    tb(slide, Inches(0.55), Inches(0.22), Inches(12), Inches(0.6), title, size=34, color=WHITE, bold=True)
    rect_bg(slide, Inches(0.55), Inches(0.8), Inches(1.0), Inches(0.035), ACCENT_BLUE)
    if subtitle:
        tb(slide, Inches(0.55), Inches(0.9), Inches(12), Inches(0.35), subtitle, size=13, color=MEDIUM_GRAY)

def gradient_accent_card(slide, l, t, w, h, color1, color2):
    """Create a card with gradient accent strip at top"""
    card(slide, l, t, w, h, CARD_BG)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l + Emu(19050), t + Emu(19050), w - Emu(38100), Inches(0.05))
    strip.fill.solid(); strip.fill.fore_color.rgb = color1; strip.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# S1  TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide); add_title_bar(slide)
rect_bg(slide, Inches(0), Inches(0.05), Inches(0.06), Inches(7.45), ACCENT_BLUE)

tb(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(1.5), 'AGENTIC DEVELOPMENT\nWITH CLAUDE CODE', size=56, color=WHITE, bold=True)
tb(slide, Inches(0.6), Inches(2.9), Inches(12), Inches(0.5), 'From AI Assistant to AI Engineering Team', size=26, color=ACCENT_BLUE)
rect_bg(slide, Inches(0.6), Inches(3.6), Inches(1.3), Inches(0.035), ACCENT_GREEN)
tb(slide, Inches(0.6), Inches(3.9), Inches(10), Inches(0.35), 'A comprehensive guide for developers mastering AI-assisted development', size=15, color=LIGHT_GRAY)
tb(slide, Inches(0.6), Inches(6.7), Inches(6), Inches(0.25), 'Claude Code + Agentic Workflow = The New SDLC', size=12, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# S2  AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, 'Agenda')

agenda_sections = [
    ('PART 1', 'Claude Code Knowledge Base', 'Skills, Commands, MCP, Sub-agents, Context, Memory, Hooks, Plugins', ACCENT_BLUE,
     ['Skills System — Custom SKILL.md files & bundled skills for repeatable workflows',
      'Commands — Built-in slash commands & bundled skill commands',
      'MCP Protocol — Connect 300+ external tools & APIs via Model Context Protocol',
      'Sub-agents — Explore, Plan, General-purpose & custom delegation',
      'Context & Memory — CLAUDE.md, auto-memory, context management tools',
      'Hooks & Plugins — Customize behavior and extend capabilities']),
    ('PART 2', 'The Agentic Development Workflow', '6-phase mandatory lifecycle with human approval gates', ACCENT_GREEN,
     ['5-Phase Operating Model — Plan, Review, Execute, Verify, Reflect',
      '6-Phase Detailed Workflow — Scratchpad through Reflect with gates',
      'Approval Gates — Human-in-the-loop at every phase (HARD STOP)',
      'Contract-First Design — API specs, event schemas, data models BEFORE code',
      'Vertical Slice Strategy — Prove 1 flow E2E before building all services',
      'Safety Boundaries — What the AI must never do — non-negotiable']),
]

y = Inches(1.15)
for label, title, subtitle, color, items in agenda_sections:
    card(slide, Inches(0.35), y, Inches(12.55), Inches(2.85))
    pill(slide, Inches(0.5), y + Inches(0.12), Inches(1.2), Inches(0.3), color, label, 10)
    tb(slide, Inches(1.9), y + Inches(0.1), Inches(9.5), Inches(0.32), title, size=22, color=WHITE, bold=True)
    tb(slide, Inches(1.9), y + Inches(0.43), Inches(9.5), Inches(0.22), subtitle, size=11, color=MEDIUM_GRAY)
    for i, item in enumerate(items):
        ix = Inches(1.9) + (i % 2) * Inches(5.2)
        iy = y + Inches(0.8) + (i // 2) * Inches(0.55)
        tb(slide, ix, iy, Inches(5.0), Inches(0.45), '  ' + item, size=11, color=LIGHT_GRAY)
    y += Inches(3.05)


# ═══════════════════════════════════════════════════════════════
# S3  WHAT IS CLAUDE CODE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, 'What is Claude Code?',
    'An agentic coding tool that reads your codebase, edits files, runs commands, and integrates with your dev tools.')

caps = [
    ('Context-Aware', ACCENT_BLUE,   'Reads entire project\nstructure, docs,\nconventions, and history'),
    ('Multi-Surface', ACCENT_GREEN,  'Terminal CLI, VS Code,\nJetBrains, Desktop,\nWeb browser, iOS'),
    ('Tool-Rich',     ACCENT_ORANGE, 'Bash, Read, Write,\nEdit, Glob, Grep,\nTask, WebFetch, MCP'),
    ('Agentic',       ACCENT_PURPLE, 'Plans, researches,\nwrites code, runs tests,\ncreates commits & PRs'),
    ('Extensible',    ACCENT_RED,    'Skills, Hooks, Plugins,\nMCP servers, custom\nsub-agents, Agent SDK'),
]
for i, (title, color, desc) in enumerate(caps):
    x = Inches(0.3 + i * 2.55)
    card(slide, x, Inches(1.45), Inches(2.4), Inches(1.85))
    pill(slide, x + Inches(0.15), Inches(1.5), Inches(2.1), Inches(0.35), color, title, 12)
    tb(slide, x + Inches(0.1), Inches(1.98), Inches(2.2), Inches(1.15), desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# How it works
card(slide, Inches(0.3), Inches(3.55), Inches(12.7), Inches(3.55))
tb(slide, Inches(0.5), Inches(3.65), Inches(10), Inches(0.28), 'How Claude Code Works', size=17, color=ACCENT_BLUE, bold=True)

flow_steps = [
    ('1. Read Context', ACCENT_BLUE, 'Loads CLAUDE.md, project\nsettings, MCP servers,\nand installed skills'),
    ('2. Understand Task', ACCENT_GREEN, 'Analyzes your request\nagainst codebase, context,\nand accumulated memory'),
    ('3. Plan Approach', ACCENT_ORANGE, 'Creates scratchpad, defines\nscope, identifies relevant\nfiles and modules'),
    ('4. Execute', ACCENT_PURPLE, 'Reads, writes, edits files,\nruns commands, spawns\nsub-agents as needed'),
    ('5. Verify & Report', ACCENT_RED, 'Runs tests, lint, type\nchecks. Reports results\nand captures learnings'),
]
for i, (title, color, desc) in enumerate(flow_steps):
    x = Inches(0.45 + i * 2.5)
    card(slide, x, Inches(4.05), Inches(2.3), Inches(2.7), CARD_HOVER, shadow=False)
    pill(slide, x + Inches(0.2), Inches(4.1), Inches(1.9), Inches(0.3), color, title, 11)
    tb(slide, x + Inches(0.1), Inches(4.55), Inches(2.1), Inches(2.0), desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.22), Inches(5.1), Inches(0.28), Inches(0.18))
        a.fill.solid(); a.fill.fore_color.rgb = BORDER_SUBTLE; a.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# S4  SKILLS SYSTEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, 'Claude Code Skills',
    'Create custom SKILL.md files to package repeatable workflows and domain knowledge.')

# LEFT — anatomy
card(slide, Inches(0.3), Inches(1.25), Inches(4.85), Inches(5.95))
tb(slide, Inches(0.5), Inches(1.35), Inches(4.0), Inches(0.28), 'Skill Anatomy', size=16, color=ACCENT_BLUE, bold=True)

mtb(slide, Inches(0.5), Inches(1.7), Inches(4.4), Inches(1.4), [
    'my-skill/',
    '  SKILL.md              (required)',
    '    YAML frontmatter + instructions',
    '',
    '  template.md           (optional)',
    '    Template for Claude to fill in',
    '',
    '  examples/             (optional)',
    '    sample.md            Expected output',
    '',
    '  scripts/              (optional)',
    '    helper.py            Executed scripts',
], size=10, color=LIGHT_GRAY, ls=1.1)

# Frontmatter tags
tb(slide, Inches(0.5), Inches(3.3), Inches(4.0), Inches(0.22), 'Frontmatter Fields', size=13, color=ACCENT_ORANGE, bold=True)
fields = [
    ('name', ACCENT_BLUE), ('description', ACCENT_BLUE), ('when_to_use', ACCENT_BLUE),
    ('arguments', ACCENT_GREEN), ('disable-model-invocation', ACCENT_GREEN), ('user-invocable', ACCENT_GREEN),
    ('allowed-tools', ACCENT_ORANGE), ('model', ACCENT_ORANGE), ('effort', ACCENT_ORANGE),
    ('context', ACCENT_PURPLE), ('agent', ACCENT_PURPLE), ('paths', ACCENT_PURPLE),
    ('hooks', ACCENT_RED), ('shell', ACCENT_RED),
]
for i, (fld, c) in enumerate(fields):
    col = i % 3; row = i // 3
    pill(slide, Inches(0.5 + col * 1.53), Inches(3.6 + row * 0.3), Inches(1.42), Inches(0.24), c, fld, 8)

# Locations
card(slide, Inches(0.3), Inches(5.55), Inches(4.85), Inches(1.65))
tb(slide, Inches(0.5), Inches(5.65), Inches(4.0), Inches(0.22), 'Where Skills Live', size=13, color=ACCENT_PURPLE, bold=True)
locs = [
    ('Enterprise', 'Managed settings — all users', 'Personal', '~/.claude/skills/'),
    ('Project', '.claude/skills/ — shared via git', 'Plugin', '<plugin>/skills/'),
]
for i, (l1, d1, l2, d2) in enumerate(locs):
    ly = Inches(5.95 + i * 0.55)
    tb(slide, Inches(0.5), ly, Inches(1.2), Inches(0.2), l1, size=10, color=ACCENT_BLUE, bold=True)
    tb(slide, Inches(1.7), ly, Inches(1.8), Inches(0.2), d1, size=9, color=LIGHT_GRAY)
    tb(slide, Inches(3.7), ly, Inches(1.0), Inches(0.2), l2, size=10, color=ACCENT_BLUE, bold=True)
    tb(slide, Inches(4.7), ly, Inches(1.6), Inches(0.2), d2, size=9, color=LIGHT_GRAY)

# RIGHT — Examples
rx = Inches(5.45); rw = Inches(7.55)

# Reference skill
card(slide, rx, Inches(1.25), rw, Inches(2.15))
tb(slide, rx + Inches(0.2), Inches(1.35), Inches(7.0), Inches(0.25), 'Reference Content Skill', size=14, color=ACCENT_GREEN, bold=True)
pill(slide, rx + Inches(0.2), Inches(1.6), Inches(7.0), Inches(0.02), ACCENT_GREEN, '', 6)
mtb(slide, rx + Inches(0.2), Inches(1.7), Inches(7.0), Inches(1.5), [
    '---',
    'description: "API design patterns for this codebase"',
    '---',
    '',
    'When writing API endpoints:',
    '  • Use RESTful naming conventions',
    '  • Return consistent error formats (RFC 7807)',
    '  • Include request validation on all inputs',
    '  • Use cursor-based pagination, not offset',
    '  • Version via URL path: /v1/resource',
], size=10, color=LIGHT_GRAY, ls=1.1)

# Task skill
card(slide, rx, Inches(3.65), rw, Inches(2.05))
tb(slide, rx + Inches(0.2), Inches(3.75), Inches(7.0), Inches(0.25), 'Task Content Skill (with /deploy example)', size=14, color=ACCENT_ORANGE, bold=True)
pill(slide, rx + Inches(0.2), Inches(4.0), Inches(7.0), Inches(0.02), ACCENT_ORANGE, '', 6)
mtb(slide, rx + Inches(0.2), Inches(4.1), Inches(7.0), Inches(1.4), [
    '---',
    'name: deploy',
    'description: Deploy the application to production',
    'disable-model-invocation: true',
    'allowed-tools: Bash(git *) Bash(npm *)',
    '---',
    'Deploy $ARGUMENTS to production:',
    '  1. Run the test suite    2. Build the application',
    '  3. Push to deployment    4. Verify deployment',
], size=10, color=LIGHT_GRAY, ls=1.1)

# Key concept cards
card(slide, rx, Inches(5.95), rw, Inches(1.25))
tb(slide, rx + Inches(0.2), Inches(6.05), Inches(7.0), Inches(0.2), 'Key Concepts', size=13, color=ACCENT_BLUE, bold=True)
concepts = [
    ('Dynamic Context', '!`cmd` injects live shell output into skill content before Claude sees it'),
    ('Fork Mode', 'context: fork runs skill in isolated sub-agent (Explore, Plan, custom)'),
    ('Invocation Control', 'disable-model-invocation restricts to manual /skill-name only'),
]
for i, (c_title, c_desc) in enumerate(concepts):
    cy = Inches(6.3 + i * 0.3)
    tb(slide, rx + Inches(0.2), cy, Inches(1.8), Inches(0.18), c_title, size=10, color=ACCENT_CYAN, bold=True)
    tb(slide, rx + Inches(2.1), cy, Inches(5.0), Inches(0.18), c_desc, size=9, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S5  COMMANDS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 5, 'Claude Code Commands',
    'Built-in slash commands and bundled skills. Type / to see everything available in-session.')

cmd_groups = [
    ('Context & Memory', ACCENT_BLUE, [
        '/compact — Summarize conversation, free up context space',
        '/clear   — Start fresh conversation (old in /resume)',
        '/memory  — Edit CLAUDE.md, manage auto-memory entries',
        '/init    — Initialize project with CLAUDE.md guide',
        '/context — Visualize current context usage as colored grid',
        '/add-dir — Add directory for file access this session',
    ]),
    ('Development', ACCENT_GREEN, [
        '/review          — Review a pull request locally',
        '/security-review — Analyze changes for vulnerabilities',
        '/batch           — Orchestrate large-scale changes in parallel',
        '/debug           — Enable debug logging, troubleshoot issues',
        '/plan            — Enter plan mode for read-only research',
        '/diff            — Interactive diff viewer of all changes',
    ]),
    ('Configuration', ACCENT_ORANGE, [
        '/model       — Switch AI models, adjust reasoning effort',
        '/permissions — Manage tool allow/ask/deny rules',
        '/config      — Open settings (theme, model, output style)',
        '/mcp         — Manage MCP server connections and OAuth',
        '/hooks       — View and configure lifecycle hook events',
        '/agents      — Create and manage sub-agent definitions',
    ]),
    ('Workflow & Automation', ACCENT_PURPLE, [
        '/schedule — Create routines on Anthropic infrastructure',
        '/loop     — Run prompt repeatedly at specified intervals',
        '/resume   — Resume a conversation by session ID or name',
        '/rewind   — Rewind conversation and/or code to a checkpoint',
        '/export   — Export current conversation as plain text',
        '/recap    — Generate a one-line summary of current session',
    ]),
]
for i, (cat, color, items) in enumerate(cmd_groups):
    col = i % 2; row = i // 2
    x = Inches(0.3 + col * 6.55)
    y = Inches(1.25 + row * 3.0)
    card(slide, x, y, Inches(6.3), Inches(2.7))
    pill(slide, x + Inches(0.12), y + Inches(0.1), Inches(2.5), Inches(0.32), color, cat, 11)
    for j, item in enumerate(items):
        iy = y + Inches(0.6 + j * 0.34)
        tb(slide, x + Inches(0.15), iy, Inches(6.0), Inches(0.3), '  ' + item, size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S6  MCP — MODEL CONTEXT PROTOCOL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 6, 'MCP — Model Context Protocol',
    'Connect Claude Code to hundreds of external tools, databases, and APIs through an open standard.')

# LEFT — What can do
card(slide, Inches(0.3), Inches(1.25), Inches(5.75), Inches(3.1))
tb(slide, Inches(0.5), Inches(1.35), Inches(5.2), Inches(0.25), 'What You Can Do with MCP', size=16, color=ACCENT_GREEN, bold=True)
mtb(slide, Inches(0.5), Inches(1.7), Inches(5.2), Inches(2.4), [
    '  Implement features from Jira issues, create PRs on GitHub',
    '  Analyze Sentry errors, query PostgreSQL databases',
    '  Upload designs from Figma, read docs from Google Drive',
    '  Create Gmail drafts, update tickets, pull Slack data',
    '  Connect 300+ servers: Notion, Asana, Stripe, Airtable',
    '  Build custom MCP servers with the MCP SDK',
    '',
    '"Connect a server when you find yourself copying data',
    'into chat from another tool or service."',
], size=11, color=LIGHT_GRAY, ls=1.25)

# RIGHT — Install + scopes
card(slide, Inches(6.35), Inches(1.25), Inches(6.65), Inches(1.45))
tb(slide, Inches(6.55), Inches(1.35), Inches(6.0), Inches(0.25), 'Installation Methods', size=15, color=ACCENT_BLUE, bold=True)
mtb(slide, Inches(6.55), Inches(1.68), Inches(6.2), Inches(0.85), [
    '  HTTP (recommended):   claude mcp add --transport http <name> <url>',
    '  Stdio (local process): claude mcp add --transport stdio <name> -- <cmd>',
    '  JSON import:           claude mcp add-json <name> \'<json>\'',
], size=10, color=LIGHT_GRAY, ls=1.3)

card(slide, Inches(6.35), Inches(2.95), Inches(6.65), Inches(1.4))
tb(slide, Inches(6.55), Inches(3.05), Inches(6.0), Inches(0.25), 'Installation Scopes', size=15, color=ACCENT_PURPLE, bold=True)
mtb(slide, Inches(6.55), Inches(3.38), Inches(6.2), Inches(0.8), [
    '  Local    — Current project only (private, stored in ~/.claude.json)',
    '  Project  — .mcp.json in project root (team-shared via version control)',
    '  User     — All your projects (private, stored in ~/.claude.json)',
], size=10, color=LIGHT_GRAY, ls=1.3)

# Features bottom
card(slide, Inches(0.3), Inches(4.65), Inches(12.7), Inches(2.55))
tb(slide, Inches(0.5), Inches(4.75), Inches(10), Inches(0.25), 'Key MCP Features', size=16, color=ACCENT_YELLOW, bold=True)

mcp_features = [
    ('Tool Search', 'Only tool names load at session start. Schemas deferred until needed — minimal context impact.'),
    ('OAuth 2.0', 'Secure authentication. Tokens stored in system keychain, refreshed automatically.'),
    ('Dynamic Updates', 'list_changed notifications update tools/prompts without disconnect/reconnect.'),
    ('MCP Resources', '@-mention server resources. Auto-fetched and included as message attachments.'),
    ('Channels', 'Servers push messages into session — react to CI results, alerts, chat messages.'),
    ('Managed MCP', 'Enterprise: managed-mcp.json for fixed server sets + allowlist/denylist policy control.'),
]
for i, (title, desc) in enumerate(mcp_features):
    col = i % 3; row = i // 3
    x = Inches(0.5 + col * 4.2); y = Inches(5.15 + row * 0.95)
    pill(slide, x, y, Inches(1.4), Inches(0.26), ACCENT_BLUE, title, 10)
    tb(slide, x, y + Inches(0.32), Inches(3.9), Inches(0.5), desc, size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S7  SUB-AGENTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 7, 'Sub-agents & Agent Delegation',
    'Spawn specialized AI assistants for specific tasks. Preserve context by isolating verbose operations.')

agents = [
    ('Explore', ACCENT_BLUE, 'Read-only research', 'Model: Haiku (fast)\nTools: Glob, Grep, Read', 'File discovery, code search,\ncodebase exploration.\nKeeps search results\nout of main context.'),
    ('Plan', ACCENT_GREEN, 'Research for planning', 'Model: Inherits\nTools: Read-only', 'Plan-mode research.\nGathers context before\npresenting an\nimplementation plan.'),
    ('General-purpose', ACCENT_ORANGE, 'Full-capability', 'Model: Inherits\nTools: All tools', 'Complex multi-step tasks\nneeding both exploration\nand modification.\nMost flexible built-in.'),
    ('Custom', ACCENT_PURPLE, 'Your own agents', 'Model: Any\nTools: You define', 'Create via /agents or\n.claude/agents/*.md.\nCustom prompt, tools,\nmodel, permissions.'),
]
for i, (name, color, desc, detail, purpose) in enumerate(agents):
    x = Inches(0.25 + i * 3.22)
    card(slide, x, Inches(1.25), Inches(3.0), Inches(3.55))
    pill(slide, x + Inches(0.2), Inches(1.3), Inches(2.6), Inches(0.4), color, name, 15)
    tb(slide, x + Inches(0.15), Inches(1.8), Inches(2.7), Inches(0.2), desc, size=10, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.1), Inches(2.05), Inches(2.8), Inches(0.45), detail, size=9, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
    pill(slide, x + Inches(0.1), Inches(2.6), Inches(2.8), Inches(0.02), color, '', 6)
    tb(slide, x + Inches(0.1), Inches(2.75), Inches(2.8), Inches(1.8), purpose, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom panels
card(slide, Inches(0.25), Inches(5.1), Inches(6.0), Inches(2.1))
tb(slide, Inches(0.45), Inches(5.2), Inches(5.5), Inches(0.22), 'Custom Agent Configuration', size=14, color=ACCENT_GREEN, bold=True)
mtb(slide, Inches(0.5), Inches(5.5), Inches(5.5), Inches(1.5), [
    'Key fields: name, description, tools, disallowedTools,',
    'model, permissionMode, maxTurns, skills, mcpServers,',
    'hooks, memory, background, effort, isolation, color,',
    'initialPrompt',
    '',
    'Locations:',
    '  .claude/agents/  — Project (shared via git)',
    '  ~/.claude/agents/ — Personal (all projects)',
    '  --agents flag    — CLI session only',
], size=9, color=LIGHT_GRAY, ls=1.2)

card(slide, Inches(6.55), Inches(5.1), Inches(6.45), Inches(2.1))
tb(slide, Inches(6.75), Inches(5.2), Inches(5.5), Inches(0.22), 'Common Delegation Patterns', size=14, color=ACCENT_ORANGE, bold=True)
mtb(slide, Inches(6.8), Inches(5.5), Inches(5.8), Inches(1.5), [
    'Isolate high-volume ops — tests, logs, file searches stay',
    'in sub-agent; only the summary returns to main context.',
    '',
    'Parallel research — spawn multiple agents exploring',
    'different modules simultaneously, then synthesize.',
    '',
    'Chain sub-agents — reviewer finds issues → fixer',
    'implements → tester verifies. Focused scope per agent.',
    '',
    'Fork mode — sub-agent inherits full conversation context,',
    'works on side task while you continue in main session.',
], size=9, color=LIGHT_GRAY, ls=1.2)


# ═══════════════════════════════════════════════════════════════
# S8  CONTEXT & MEMORY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 8, 'Context & Memory Management',
    'How Claude Code understands your project and remembers across sessions.')

# LEFT — CLAUDE.md
card(slide, Inches(0.3), Inches(1.25), Inches(6.25), Inches(3.7))
tb(slide, Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.25), 'CLAUDE.md — Your Project Constitution', size=16, color=ACCENT_BLUE, bold=True)
mtb(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(3.0), [
    'The primary way to teach Claude about your project.',
    'Claude reads CLAUDE.md at the start of EVERY session.',
    '',
    'What to include:',
    '  • Coding standards and architecture decisions',
    '  • Preferred libraries, frameworks, and tools',
    '  • Common commands (build, test, lint, deploy)',
    '  • Key patterns, naming conventions, import rules',
    '  • Review checklists and security requirements',
    '  • Path-specific rules for monorepos',
    '',
    'Locations (merged in priority order):',
    '  • ./CLAUDE.md         — Project root (check into git)',
    '  • .claude/CLAUDE.md   — Same as root, alternate location',
    '  • ~/.claude/CLAUDE.md — Personal (all your projects)',
    '  • Managed settings    — Enterprise-wide policies',
], size=10, color=LIGHT_GRAY, ls=1.12)

# RIGHT — auto memory + context
card(slide, Inches(6.85), Inches(1.25), Inches(6.15), Inches(1.7))
tb(slide, Inches(7.05), Inches(1.35), Inches(5.5), Inches(0.22), 'Auto Memory', size=16, color=ACCENT_GREEN, bold=True)
mtb(slide, Inches(7.1), Inches(1.65), Inches(5.6), Inches(1.1), [
    'Claude builds auto-memory as it works — saving learnings',
    'like build commands, debugging insights, and project',
    'patterns across sessions without you writing anything.',
    '',
    'Manage with /memory. Enable/disable per project.',
    'Auto-memory entries survive conversation compaction.',
], size=10, color=LIGHT_GRAY, ls=1.15)

card(slide, Inches(6.85), Inches(3.15), Inches(6.15), Inches(1.8))
tb(slide, Inches(7.05), Inches(3.25), Inches(5.5), Inches(0.22), 'Context Management Tools', size=16, color=ACCENT_ORANGE, bold=True)
mtb(slide, Inches(7.1), Inches(3.55), Inches(5.6), Inches(1.2), [
    '/compact  — Summarize conversation, free context space',
    '/context  — Visualize context usage as colored grid',
    '/add-dir  — Add directory for file access this session',
    '/clear    — Start fresh (old session in /resume)',
    '',
    'Auto-compaction triggers at ~95% context capacity.',
    'Preserves rules, skills, and recent messages across',
    'compaction boundaries for continuity.',
], size=10, color=LIGHT_GRAY, ls=1.15)

# Architecture bottom
card(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(2.0))
tb(slide, Inches(0.5), Inches(5.3), Inches(10), Inches(0.22), 'Memory & Context Architecture', size=15, color=ACCENT_YELLOW, bold=True)

arch = [
    ('System Prompt', ACCENT_BLUE, 'Base instructions every session — persona, safety rules, tool usage.'),
    ('CLAUDE.md Files', ACCENT_GREEN, 'Project root, user home, managed — merged in priority order at session start.'),
    ('Auto Memory', ACCENT_ORANGE, 'Cross-session learnings: build commands, debugging insights, project patterns.'),
    ('Sub-agent Memory', ACCENT_PURPLE, 'Custom agents get persistent memory dirs (user/project/local) for learning.'),
    ('Skills & Commands', ACCENT_RED, 'Inspected at startup, content loaded only on invocation. Near-zero idle cost.'),
    ('MCP Tools', ACCENT_BLUE, 'Deferred via Tool Search — only actually-used tools enter the context window.'),
]
for i, (title, color, desc) in enumerate(arch):
    y = Inches(5.6 + i * 0.25)
    pill(slide, Inches(0.5), y, Inches(1.8), Inches(0.2), color, title, 9)
    tb(slide, Inches(2.45), y, Inches(10.0), Inches(0.2), desc, size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S9  HOOKS, PLUGINS & AUTOMATION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 9, 'Hooks, Plugins & Automation',
    'Customize Claude Code behavior, extend capabilities, and automate recurring work.')

cols = [
    ('Hooks', ACCENT_BLUE, [
        'Lifecycle Automation',
        '',
        'PreToolUse — Validate before run',
        'PostToolUse — React after execution',
        'Notification — Desktop alerts',
        'SubagentStart / SubagentStop',
        'Elicitation — Handle MCP dialogs',
        '',
        'Examples:',
        '• Auto-format after every file edit',
        '• Run linter before git commits',
        '• Block destructive operations',
        '• Validate SQL queries as read-only',
        '• Send Slack notification on deploy',
    ]),
    ('Plugins', ACCENT_GREEN, [
        'Share & Distribute Extensions',
        '',
        'A plugin can bundle all of:',
        '  • Skills (SKILL.md files)',
        '  • MCP servers',
        '  • Sub-agents',
        '  • Hooks',
        '  • CLAUDE.md fragments',
        '',
        'Install via /plugin command.',
        'Manage: enable, disable, update.',
        'Team: share via managed settings.',
        'Build custom plugins for your org.',
    ]),
    ('Automation', ACCENT_ORANGE, [
        'Run Claude on a Schedule',
        '',
        'Routines:',
        '  • Run on Anthropic infrastructure',
        '  • Trigger: schedule, API, GitHub',
        '  • Continue even when your PC is off',
        '',
        'Desktop Scheduled Tasks:',
        '  • Run on your local machine',
        '  • Direct filesystem access',
        '',
        '/loop  — Repeat within CLI session',
        '/schedule — Create & manage routines',
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.3 + i * 4.35)
    card(slide, x, Inches(1.25), Inches(4.1), Inches(5.2))
    pill(slide, x + Inches(0.3), Inches(1.3), Inches(3.5), Inches(0.38), color, title, 16)
    mtb(slide, x + Inches(0.2), Inches(1.85), Inches(3.7), Inches(4.4), items, size=10, color=LIGHT_GRAY, ls=1.15)

# Bottom bar
card(slide, Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.55))
tb(slide, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.2),
    'Also: GitHub Actions CI/CD  |  Slack @Claude Integration  |  Chrome Browser Extension  |  Remote Control (phone/browser)  |  Agent SDK  |  Agent Teams',
    size=9, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# S10  BEST PRACTICES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 10, 'Best Practices for Claude Code',
    'Patterns that maximize effectiveness and minimize friction.')

practices = [
    ('Write Effective CLAUDE.md', ACCENT_BLUE, [
        'Keep it concise — every line costs context tokens every session',
        'State facts, not narratives: "We use PostgreSQL" not backstory',
        'Include common commands: build, test, lint, deploy, migrate',
        'Use path-specific rules for monorepos: **/frontend/** scoped',
    ]),
    ('Craft Good Prompts', ACCENT_GREEN, [
        'Be specific: "Add validation to POST /users" not "improve API"',
        'Provide context: reference files, error messages, expected behavior',
        'Break large tasks into smaller, focused, independently-verifiable requests',
        'Use @ mentions to reference files and sub-agents for precision',
    ]),
    ('Manage Context Wisely', ACCENT_ORANGE, [
        '/compact when conversation grows — preserves key information',
        'Delegate verbose operations to sub-agents (test runs, log analysis)',
        '/context visualizes what is consuming your context window',
        '/clear when context is stale; old sessions available via /resume',
    ]),
    ('Leverage Sub-agents', ACCENT_PURPLE, [
        'Create custom sub-agents for recurring task types in your project',
        'Explore agent for codebase research — isolates search from main context',
        'Run independent research in parallel with multiple sub-agents',
        'Chain: reviewer finds issues → fixer implements → tester verifies',
    ]),
    ('Build in the Open', ACCENT_RED, [
        'Commit early and often — Claude works best with git history available',
        'Use /review for pre-merge code review, /security-review for vulns',
        'Share CLAUDE.md and skills via version control for team consistency',
        'Project-scoped MCP (.mcp.json) for shared team tool access',
    ]),
]
for i, (title, color, items) in enumerate(practices):
    y = Inches(1.2 + i * 1.25)
    card(slide, Inches(0.3), y, Inches(12.7), Inches(1.05))
    pill(slide, Inches(0.45), y + Inches(0.1), Inches(3.3), Inches(0.25), color, title, 11)
    for j, item in enumerate(items):
        ix = Inches(0.5 + (j % 2) * 6.2)
        iy = y + Inches(0.45 + (j // 2) * 0.26)
        tb(slide, ix, iy, Inches(6.0), Inches(0.22), '  ' + item, size=10, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S11  SECTION BREAK — PART 2
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide); add_title_bar(slide)
rect_bg(slide, Inches(0), Inches(0.05), Inches(0.06), Inches(7.45), ACCENT_GREEN)

tb(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.7), 'PART 2', size=24, color=ACCENT_GREEN, bold=True)
tb(slide, Inches(0.6), Inches(2.9), Inches(12), Inches(1.0), 'The Agentic Development\nWorkflow', size=52, color=WHITE, bold=True)
tb(slide, Inches(0.6), Inches(4.2), Inches(10), Inches(0.35), 'Think first. Plan second. Code last. Always.', size=20, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S12  ADLC OVERVIEW — TRADITIONAL vs AGENTIC
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 12, 'The Agentic Development Lifecycle',
    'Traditional SDLC replaced by systematic AI-human collaboration with mandatory approval gates.')

# Traditional
pill(slide, Inches(0.5), Inches(1.25), Inches(5.7), Inches(0.42), ACCENT_RED, 'TRADITIONAL SDLC', 16)
card(slide, Inches(0.3), Inches(1.75), Inches(6.1), Inches(5.0))
mtb(slide, Inches(0.5), Inches(1.9), Inches(5.7), Inches(4.5), [
    'Plan → Design → Code → Test → Deploy → Maintain',
    '',
    '  • Linear, sequential phases',
    '  • Long feedback loops (weeks or months)',
    '  • Human writes all code from scratch',
    '  • Reviews happen AFTER code is written',
    '  • Testing is a late-phase activity',
    '  • Documentation often stale or missing',
    '  • Single developer mindset',
    '    "I build everything myself"',
    '',
    'Result:',
    '  Slow delivery, error-prone,',
    '  hard to course-correct mid-project',
], size=11, color=LIGHT_GRAY, ls=1.15)

# Agentic
pill(slide, Inches(6.9), Inches(1.25), Inches(5.7), Inches(0.42), ACCENT_GREEN, 'AGENTIC DEVELOPMENT LIFECYCLE', 16)
card(slide, Inches(6.8), Inches(1.75), Inches(6.1), Inches(5.0))
mtb(slide, Inches(7.0), Inches(1.9), Inches(5.7), Inches(4.5), [
    'Scratchpad → Plan → Tasks → Execute → Verify → Reflect',
    '',
    '  • Cyclic, iterative process',
    '  • Fast feedback (minutes or hours)',
    '  • AI generates code; human reviews & guides',
    '  • Architecture review BEFORE any code',
    '  • Tests written alongside code, always verified',
    '  • Documentation-first: contracts define all',
    '  • Multi-agent team',
    '    6+ specialized AI personas with boundaries',
    '',
    'Result:',
    '  Fast delivery, quality-gated,',
    '  continuously self-improving',
], size=11, color=LIGHT_GRAY, ls=1.15)

# Arrow
a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.15), Inches(3.4), Inches(0.6), Inches(0.4))
a.fill.solid(); a.fill.fore_color.rgb = ACCENT_ORANGE; a.line.fill.background()
tb(slide, Inches(5.7), Inches(3.1), Inches(1.5), Inches(0.25), 'TRANSFORM', size=9, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

tb(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.28),
    '"AI is a collaborator, not a tool. Humans set direction; agents execute."', size=12, color=ACCENT_YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# S13  5-PHASE OPERATING MODEL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 13, 'The 5-Phase Operating Model',
    'Every non-trivial task follows this lifecycle. Human approval gates enforce quality and safety.')

phases = [
    ('PLAN', ACCENT_BLUE,   'Create scratchpad\nAnalyze codebase\nDefine scope\nIdentify invariants', True),
    ('REVIEW', ACCENT_ORANGE, 'HARD STOP\nHuman reviews\nExplicit APPROVE\nrequired to continue', True),
    ('EXECUTE', ACCENT_GREEN, 'Implement the plan\nFollow invariants\nOne task at a time\nNo scope creep', True),
    ('VERIFY', ACCENT_PURPLE, 'Run tests\nLint & type check\nConfirm behavior\nmatches objective', False),
    ('REFLECT', ACCENT_YELLOW, 'Review execution\nCapture learnings\nSuggest workflow\nimprovements', False),
]
for i, (phase, color, desc, is_gate) in enumerate(phases):
    x = Inches(0.25 + i * 2.55)
    card(slide, x, Inches(1.2), Inches(2.4), Inches(0.85), color, shadow=False, border=False)
    tb(slide, x + Inches(0.1), Inches(1.27), Inches(2.2), Inches(0.38), phase, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    card(slide, x, Inches(2.15), Inches(2.4), Inches(1.9))
    tb(slide, x + Inches(0.1), Inches(2.2), Inches(2.2), Inches(1.8), desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow_r(slide, x + Inches(2.36), Inches(1.5), Inches(0.22), Inches(0.14))

# Golden rules
card(slide, Inches(0.25), Inches(4.35), Inches(12.8), Inches(2.85), ERROR_BG, border=False)
tb(slide, Inches(0.5), Inches(4.45), Inches(10), Inches(0.28), 'GOLDEN RULES', size=18, color=ACCENT_RED, bold=True)
mtb(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.1), [
    '1. Scratchpad is the SOURCE OF TRUTH. Plans explain HOW. Tasks define WHAT. Code is the LAST step.',
    '2. NEVER skip approval gates — each one is a HARD STOP requiring explicit human "APPROVE".',
    '3. NEVER generate SCRATCHPAD + PLAN in the same response. NEVER PLAN + TASKS together.',
    '4. NEVER write implementation code before TASKS are approved by a human.',
    '5. If any decision changes during execution: update scratchpad FIRST and request re-approval.',
    '6. Before coding, ALWAYS ask: Approved scratchpad? Approved plan? Approved tasks? If NO → STOP.',
], size=12, color=LIGHT_GRAY, ls=1.25)


# ═══════════════════════════════════════════════════════════════
# S14  6-PHASE MANDATORY WORKFLOW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 14, 'The 6-Phase Mandatory Workflow',
    'Every non-trivial task: Scratchpad → Plan → Tasks → Execute → Verify → Reflect.')

wf = [
    ('Phase 1', 'SCRATCHPAD', ACCENT_BLUE,   'Analyze codebase.\nUnderstand problem.\nDefine scope & invariants.\nOutput: SCRATCHPAD.md', '🛑 GATE 1'),
    ('Phase 2', 'PLAN',       ACCENT_PURPLE, 'Architecture changes.\nRisk analysis.\nEdge cases, non-goals.\nOutput: PLAN.md', '🛑 GATE 2'),
    ('Phase 3', 'TASKS',      ACCENT_ORANGE, 'Ordered steps.\nFiles to modify.\nDependencies & outputs.\nOutput: TASKS.md', '🛑 GATE 3'),
    ('Phase 4', 'IMPLEMENT',  ACCENT_GREEN,  'Write code per tasks.\nOne task at a time.\nNo scope creep.\nIf change needed, re-approve.', 'Follow tasks'),
    ('Phase 5', 'VERIFY',     ACCENT_CYAN,   'Run tests, lint,\ntype checks.\nVerify vs scratchpad.\nGenerate missing tests.', 'Quality gate'),
    ('Phase 6', 'REFLECT',    ACCENT_RED,    'Review execution.\nCapture learnings.\nSuggest improvements.\nOutput: reflection.md', 'For non-trivial'),
]
for i, (label, name, color, desc, gate) in enumerate(wf):
    x = Inches(0.15 + i * 2.18)
    card(slide, x, Inches(1.2), Inches(2.0), Inches(0.52), color, shadow=False, border=False)
    tb(slide, x + Inches(0.05), Inches(1.2), Inches(1.9), Inches(0.18), label, size=8, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.05), Inches(1.35), Inches(1.9), Inches(0.3), name, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    card(slide, x, Inches(1.82), Inches(2.0), Inches(2.2))
    tb(slide, x + Inches(0.05), Inches(1.88), Inches(1.9), Inches(2.0), desc, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    c = ACCENT_ORANGE if 'GATE' in gate else ACCENT_GREEN
    tb(slide, x + Inches(0.05), Inches(4.1), Inches(1.9), Inches(0.3), gate, size=9, color=c, bold=True, align=PP_ALIGN.CENTER)
    if i < 5: arrow_r(slide, x + Inches(1.97), Inches(2.4), Inches(0.18), Inches(0.12))

# Context loading
card(slide, Inches(0.15), Inches(4.65), Inches(13.0), Inches(2.55))
tb(slide, Inches(0.35), Inches(4.75), Inches(11), Inches(0.25), 'Context Loading — Before Every Task', size=16, color=ACCENT_BLUE, bold=True)
mtb(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.9), [
    '',
    '  PROJECT.md      — What: tech stack, commands, directory structure, architecture patterns, deployment',
    '  CONVENTIONS.md  — How: naming conventions, import ordering, code organization, type safety, error handling, tests',
    '  BOUNDARIES.md   — Don\'t: forbidden actions, scope rules, sensitive data protection, code integrity, approval gates',
    '',
    'These 3 pillars are the primary source of truth. They define what the AI can do, how it should do it, and what it must never do.',
], size=10, color=LIGHT_GRAY, ls=1.15)


# ═══════════════════════════════════════════════════════════════
# S15  PHASES 1-3 DEEP DIVE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 15, 'Phases 1-3: Think Before You Build',
    'Scratchpad, Plan, and Tasks — the "no code" phases where all thinking happens.')

p1_3 = [
    ('PHASE 1: SCRATCHPAD', ACCENT_BLUE, [
        'NO code or plans in this phase whatsoever.',
        'Read PROJECT.md, CONVENTIONS.md, BOUNDARIES.md',
        'Analyze the codebase relevant to the task',
        'Understand the problem deeply — gather all context',
        'Define scope, architecture invariants, known risks',
        'Identify open questions and assumptions',
        'Output: SCRATCHPAD.md → Human reviews → APPROVE',
    ]),
    ('PHASE 2: PLAN', ACCENT_PURPLE, [
        'Scratchpad must be APPROVED before starting.',
        'NO code or task lists in this phase.',
        'Define architecture changes, modules/files affected',
        'Risk analysis, edge cases, validation strategy',
        'Explicit non-goals — what we are NOT doing',
        'If plan contradicts scratchpad → STOP, update scratchpad',
        'Output: PLAN.md → Human reviews → APPROVE',
    ]),
    ('PHASE 3: TASKS', ACCENT_ORANGE, [
        'Plan must be APPROVED before starting.',
        'NO code in this phase — still thinking!',
        'Break plan into ordered implementation steps',
        'Files to create, modify, or delete per step',
        'Dependencies between steps clearly marked',
        'Expected output defined for each step',
        'Output: TASKS.md → Human reviews → APPROVE',
    ]),
]
for i, (title, color, items) in enumerate(p1_3):
    x = Inches(0.25 + i * 4.25)
    card(slide, x, Inches(1.2), Inches(4.0), Inches(0.5), color, shadow=False, border=False)
    tb(slide, x + Inches(0.1), Inches(1.22), Inches(3.8), Inches(0.44), title, size=16, color=WHITE, bold=True)
    card(slide, x, Inches(1.8), Inches(4.0), Inches(3.7))
    mtb(slide, x + Inches(0.15), Inches(1.9), Inches(3.7), Inches(3.4), items, size=11, color=LIGHT_GRAY, ls=1.3)
    # Gate
    pill(slide, x + Inches(0.5), Inches(5.7), Inches(3.0), Inches(0.45), ERROR_BG, '🛑 HARD STOP — GATE ' + str(i+1) + '   WAIT for APPROVE', 9, ACCENT_RED)

# Dev checklist
card(slide, Inches(0.25), Inches(6.4), Inches(12.8), Inches(0.75))
tb(slide, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.25),
    'Dev Checklist: Before coding, ask — Scratchpad approved? Plan approved? Tasks approved? If any answer is NO → STOP. Do not write any code.',
    size=12, color=ACCENT_YELLOW, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# S16  PHASES 4-6 DEEP DIVE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 16, 'Phases 4-6: Execute, Verify, Learn',
    'Implementation, testing, and reflection — where code gets written and quality is proven.')

p4_6 = [
    ('PHASE 4: IMPLEMENT', ACCENT_GREEN, [
        'Tasks must be APPROVED before writing any code',
        'Implement one task at a time in defined order',
        'Follow scratchpad invariants strictly — no deviation',
        'Match existing code style and conventions exactly',
        'If decision changes → update scratchpad FIRST → re-approve',
        'No scope creep. No unrelated refactoring.',
        'Files written ONCE per task — re-read before re-writing',
    ]),
    ('PHASE 5: VERIFY', ACCENT_CYAN, [
        'Run all relevant tests for the changes',
        'Run lint and type checks — must pass cleanly',
        'Verify behavior matches scratchpad objective',
        'Generate or update tests as needed',
        'If something deviates: explain why, propose update',
        'Check logs and metrics if applicable',
        'Tests verify BEHAVIOR, not implementation details',
    ]),
    ('PHASE 6: REFLECT', ACCENT_RED, [
        'Review execution quality — did it match the plan?',
        'Note any unexpected complexities or deviations',
        'Capture learnings for future reference and reuse',
        'Suggest workflow or convention improvements',
        'Mandatory for multi-file & architectural changes',
        'Optional for trivial tasks (single file, simple fix)',
        'Output: reflection.md with insights & suggestions',
    ]),
]
for i, (title, color, items) in enumerate(p4_6):
    x = Inches(0.25 + i * 4.25)
    card(slide, x, Inches(1.2), Inches(4.0), Inches(0.5), color, shadow=False, border=False)
    tb(slide, x + Inches(0.1), Inches(1.22), Inches(3.8), Inches(0.44), title, size=16, color=WHITE, bold=True)
    card(slide, x, Inches(1.8), Inches(4.0), Inches(5.4))
    mtb(slide, x + Inches(0.15), Inches(1.9), Inches(3.7), Inches(5.0), items, size=11, color=LIGHT_GRAY, ls=1.3)


# ═══════════════════════════════════════════════════════════════
# S17  APPROVAL GATES DETAIL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 17, 'Approval Gates: Human-in-the-Loop',
    'Four mandatory checkpoints where the AI MUST stop and wait for explicit human approval.')

gates = [
    ('GATE 1', ACCENT_BLUE,   'After SCRATCHPAD', 'Review scope, invariants, risks.\nConfirm AI understands correctly\nbefore any planning begins.',
     'SCRATCHPAD.md', '"Review scratchpad.\nReply APPROVE to\ncontinue."'),
    ('GATE 2', ACCENT_PURPLE, 'After PLAN', 'Review architecture decisions,\nmodules affected, risk analysis.\nConfirm approach before tasks.',
     'PLAN.md', '"Review the plan.\nReply APPROVE to\ncontinue."'),
    ('GATE 3', ACCENT_ORANGE, 'After TASKS', 'Review ordered implementation\nsteps. Confirm each task is clear,\nscoped, with expected outputs.',
     'TASKS.md', '"Review the tasks.\nReply APPROVE to\ncontinue."'),
    ('GATE 4', ACCENT_GREEN,  'After IMPLEMENT', 'Review completed code. Verify\nmatches plan and passes quality\nchecks before final testing.',
     'Completed Code', '"Implementation complete.\nPlease review before\nfinal testing."'),
]
for i, (gate, color, phase, desc, output, say) in enumerate(gates):
    y = Inches(1.25 + i * 1.55)
    card(slide, Inches(0.3), y, Inches(12.7), Inches(1.35))
    pill(slide, Inches(0.45), y + Inches(0.1), Inches(1.3), Inches(0.28), color, gate, 12)
    tb(slide, Inches(0.5), y + Inches(0.45), Inches(1.8), Inches(0.2), 'Checks: ' + phase, size=9, color=MEDIUM_GRAY)
    tb(slide, Inches(2.5), y + Inches(0.15), Inches(5.0), Inches(1.0), desc, size=12, color=LIGHT_GRAY)
    tb(slide, Inches(8.0), y + Inches(0.2), Inches(2.0), Inches(0.6), 'Artifact:\n' + output, size=11, color=ACCENT_YELLOW, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, Inches(10.2), y + Inches(0.15), Inches(2.5), Inches(0.9), say, size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

tb(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.2),
    'Additional triggers: Architecture decisions | Breaking changes | New dependencies | Plan changes mid-execution → ALL require re-approval.',
    size=9, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# S18  CONTRACT-FIRST DEVELOPMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 18, 'Contract-First Development',
    'Design ALL contracts before writing ANY code. Contracts are the source of truth — code conforms to them.')

contracts = [
    ('API Contracts', ACCENT_BLUE, [
        'OpenAPI specs for EVERY endpoint',
        'kebab-case paths, plural collections',
        'RFC 7807 standard error responses',
        'Cursor-based pagination (not offset)',
        'URL-path versioning (/v1/)',
        'Rate limiting per-client tier',
    ]),
    ('Event Contracts', ACCENT_GREEN, [
        'Standard event envelope schema:',
        '{type, schemaVersion, source,',
        'correlationId, timestamp, payload}',
        'Topic catalog with partition keys',
        'Outbox pattern for transactional publishing',
        'Inbox pattern for deduplication (idempotent)',
        'Schema evolution: add=OK, remove=major, rename=never',
    ]),
    ('Data Contracts', ACCENT_ORANGE, [
        'ER diagrams with index strategy included',
        'Storage type per service (PG, Redis, S3)',
        'Connection pooling strategy (PgBouncer / proxy)',
        '1yr and 3yr storage growth projections',
        'Database-per-service — no cross-service JOINs',
        'Numbered migrations, never edit existing files',
    ]),
]
for i, (title, color, items) in enumerate(contracts):
    x = Inches(0.25 + i * 4.3)
    card(slide, x, Inches(1.2), Inches(4.05), Inches(0.45), color, shadow=False, border=False)
    tb(slide, x + Inches(0.1), Inches(1.22), Inches(3.85), Inches(0.4), title, size=17, color=WHITE, bold=True)
    card(slide, x, Inches(1.75), Inches(4.05), Inches(4.1))
    mtb(slide, x + Inches(0.15), Inches(1.85), Inches(3.75), Inches(3.8), ['  ' + it for it in items], size=11, color=LIGHT_GRAY, ls=1.25)

# Bottom
card(slide, Inches(0.25), Inches(6.15), Inches(7.4), Inches(1.05), SUCCESS_BG, border=False)
mtb(slide, Inches(0.45), Inches(6.2), Inches(7.0), Inches(0.9), [
    '  BENEFITS: Prevents #1 production issue (inconsistent APIs)',
    '  Enables parallel development | Contracts serve as living documentation',
    '  API consumers can start integration before backend is implemented',
], size=10, color=LIGHT_GRAY, ls=1.2)

card(slide, Inches(7.95), Inches(6.15), Inches(5.05), Inches(1.05), ERROR_BG, border=False)
mtb(slide, Inches(8.1), Inches(6.2), Inches(4.7), Inches(0.9), [
    '  COMMON MISTAKES:',
    '  Code-first APIs → inconsistent contracts',
    '  Offset pagination → degrades at page 1000+',
    '  No event versioning → silently breaking consumers',
], size=10, color=LIGHT_GRAY, ls=1.2)


# ═══════════════════════════════════════════════════════════════
# S19  VERTICAL SLICE STRATEGY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 19, 'Vertical Slice Strategy',
    'Prove ONE complete flow works end-to-end before building all services horizontally.')

# DON'T
pill(slide, Inches(0.5), Inches(1.25), Inches(5.5), Inches(0.42), ACCENT_RED, 'HORIZONTAL BUILD  (DON\'T)', 16)
card(slide, Inches(0.3), Inches(1.75), Inches(5.9), Inches(2.15))
mtb(slide, Inches(0.5), Inches(1.9), Inches(5.5), Inches(1.8), [
    '██  All APIs built first',
    '██  All Services built first',
    '██  All Databases built',
    '██  All Events wired up',
    '',
    'Nothing works end-to-end.',
    'Integration issues discovered',
    'only when everything is wired.',
    'Late discoveries = expensive fixes.',
], size=11, color=LIGHT_GRAY, ls=1.2)

# DO
pill(slide, Inches(6.9), Inches(1.25), Inches(5.5), Inches(0.42), ACCENT_GREEN, 'VERTICAL SLICE  (DO THIS)', 16)
card(slide, Inches(6.8), Inches(1.75), Inches(5.9), Inches(2.15))
mtb(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(1.8), [
    '██  One API endpoint',
    '██  One Service logic',
    '██  One Database table',
    '██  One Event published',
    '',
    'Complete flow works end-to-end.',
    'Architecture validated early.',
    'Problems found cheaply,',
    'before scaling to all services.',
], size=11, color=LIGHT_GRAY, ls=1.2)

# Arrow
a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.05), Inches(2.5), Inches(0.65), Inches(0.35))
a.fill.solid(); a.fill.fore_color.rgb = ACCENT_GREEN; a.line.fill.background()

# Build tiers
tb(slide, Inches(0.5), Inches(4.2), Inches(11), Inches(0.35), 'After Vertical Slice Success — Build Tier-by-Tier', size=18, color=WHITE, bold=True)

tiers = [
    ('Tier 1: Foundation', ACCENT_GREEN,  'User/Auth, Config, API Gateway — shared infrastructure'),
    ('Tier 2: Core Business', ACCENT_BLUE, 'Product/Catalog, Inventory, Pricing — domain logic'),
    ('Tier 3: Transactions', ACCENT_ORANGE, 'Order, Payment, Cart — revenue-critical flows'),
    ('Tier 4: Support Systems', ACCENT_PURPLE, 'Notifications, Analytics, Admin — operational needs'),
]
for i, (tier, color, examples) in enumerate(tiers):
    y = Inches(4.75 + i * 0.65)
    pill(slide, Inches(0.5), y, Inches(2.8), Inches(0.45), color, tier, 12)
    tb(slide, Inches(3.5), y + Inches(0.08), Inches(9.0), Inches(0.32), examples, size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S20  SAFETY BOUNDARIES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 20, 'Safety Boundaries: What the AI Must Never Do',
    'Hard constraints enforced across all tools and workflows. Non-negotiable.')

boundaries = [
    ('Scope Control', ACCENT_ORANGE, [
        'Do NOT modify code outside the task scope',
        'Do NOT refactor unrelated code while fixing bugs',
        'If > 3 files need changes, verify scope first',
        'Each file written once — re-read before re-writing',
    ]),
    ('Destructive Operations', ACCENT_RED, [
        'Never rm -rf or equivalent destructive commands',
        'Never git push or git reset --hard without approval',
        'Never modify build outputs or node_modules directly',
        'Never modify dependency lock files directly',
    ]),
    ('Sensitive Data', ACCENT_RED, [
        'Never read .env or environment variable files',
        'Never read secrets, credentials, or API keys',
        'Never read Terraform state or infrastructure secrets',
        'Never expose secrets in code, logs, or documentation',
    ]),
    ('Code Integrity', ACCENT_PURPLE, [
        'Never mix editing strategies on the same file',
        'Never wipe entire files — edit section by section',
        'If file changes externally, re-read before writing',
        'Read all related files first, then write one by one',
    ]),
    ('Approval Gates', ACCENT_RED, [
        'Never skip a human approval gate — CRITICAL violation',
        'Never generate SCRATCHPAD + PLAN in same response',
        'Never generate PLAN + TASKS in same response',
        'Never write code before TASKS are approved by human',
    ]),
    ('Tool Discipline', ACCENT_BLUE, [
        'Doc tools: for external library documentation only',
        'Analysis tools: symbols only — not broad text search',
        'Repo tools: PR management, code review, issues only',
        'Never commit secrets through any tool or integration',
    ]),
]
for i, (title, color, items) in enumerate(boundaries):
    row = i // 3; col = i % 3
    x = Inches(0.25 + col * 4.3); y = Inches(1.25 + row * 3.0)
    card(slide, x, y, Inches(4.05), Inches(2.7))
    pill(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.2), Inches(0.3), color, title, 10)
    mtb(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.75), Inches(2.0), ['  ' + it for it in items], size=10, color=LIGHT_GRAY, ls=1.25)


# ═══════════════════════════════════════════════════════════════
# S21  COMMON MISTAKES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 21, 'Common Mistakes & Anti-patterns',
    'What to avoid at each stage of the agentic development workflow.')

mistakes = [
    ('Scratchpad Phase', ACCENT_BLUE,   'Jumping into code without analysis',
     'Skipping context file review entirely', 'Not defining explicit invariants and scope'),
    ('Plan Phase',       ACCENT_PURPLE, 'Vague plans without specific files',
     'No edge case or risk analysis done', 'Plans that ignore scratchpad invariants'),
    ('Tasks Phase',      ACCENT_ORANGE, 'Tasks too large to verify independently',
     'Missing dependencies between task steps', 'No expected output defined per task'),
    ('Execute Phase',    ACCENT_GREEN,  'Scope creep — adding unapproved changes',
     'Mixing refactoring with feature work', 'Ignoring existing code style conventions'),
    ('Verify Phase',     ACCENT_CYAN,   'Skipping tests because "it works locally"',
     'Not checking for regressions in related areas', 'Ignoring lint and type check failures'),
    ('Reflect Phase',    ACCENT_RED,    'Skipping reflection entirely after execution',
     'Not capturing learnings for future sessions', 'Repeating the same mistakes across tasks'),
]
for i, (phase, color, m1, m2, m3) in enumerate(mistakes):
    y = Inches(1.25 + i * 1.02)
    card(slide, Inches(0.25), y, Inches(12.8), Inches(0.82))
    pill(slide, Inches(0.4), y + Inches(0.12), Inches(1.9), Inches(0.25), color, phase, 10)
    tb(slide, Inches(2.5), y + Inches(0.12), Inches(3.2), Inches(0.25), '✗ ' + m1, size=10, color=ACCENT_RED)
    tb(slide, Inches(5.9), y + Inches(0.12), Inches(3.2), Inches(0.25), '✗ ' + m2, size=10, color=ACCENT_RED)
    tb(slide, Inches(9.3), y + Inches(0.12), Inches(3.2), Inches(0.25), '✗ ' + m3, size=10, color=ACCENT_RED)
    tb(slide, Inches(0.4), y + Inches(0.5), Inches(12.0), Inches(0.22),
        'Correct approach: Follow phase rules. Approval gates exist to prevent exactly these mistakes.',
        size=9, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# S22  GETTING STARTED
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 22, 'Getting Started with Agentic Development',
    'Your first steps into AI-assisted software engineering.')

steps = [
    ('1', 'Set Up Claude Code', 'Install via curl/winget/brew. Run claude login. Explore the CLI and VS Code extension capabilities.'),
    ('2', 'Create Your CLAUDE.md', 'Write your project constitution: tech stack, coding standards, common commands, and architecture patterns.'),
    ('3', 'Learn the Workflow Rhythm', 'For any task: scratchpad first (think), plan second (design), tasks third (break down). Get approval at each gate.'),
    ('4', 'Build Reusable Skills', 'Create .claude/skills/ for repeatable workflows: /deploy, /review-pr, /commit. Share with team via version control.'),
    ('5', 'Connect MCP Servers', 'Connect to GitHub, Jira, databases, and monitoring tools. Let Claude read and act on your existing tools directly.'),
    ('6', 'Create Custom Sub-agents', 'Define agents for recurring tasks: code-reviewer, test-runner, security-auditor. Use the /agents command interface.'),
    ('7', 'Automate & Continuously Evolve', 'Set up CI/CD integration. Create routines for scheduled tasks. Review and refine your setup regularly.'),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.2 + i * 0.88)
    card(slide, Inches(0.3), y, Inches(12.7), Inches(0.68))
    pill(slide, Inches(0.45), y + Inches(0.08), Inches(0.42), Inches(0.42), ACCENT_BLUE, num, 16)
    tb(slide, Inches(1.05), y + Inches(0.1), Inches(3.5), Inches(0.26), title, size=16, color=WHITE, bold=True)
    tb(slide, Inches(1.05), y + Inches(0.38), Inches(11.5), Inches(0.25), desc, size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S23  MENTAL MODELS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 23, 'Mental Models for Agentic Development',
    'How to think when collaborating with AI agents.')

models = [
    ('The Operating Model', 'Plan → Review → Execute → Verify → Reflect', ACCENT_BLUE,
     'Before coding, ask: Approved scratchpad? Approved plan? Approved tasks? If any answer is NO — STOP. Code is always the last step.'),
    ('Discovery Triangle', 'WHAT (Problem) ← WHO (Users) → WHY (Value) → Risks', ACCENT_PURPLE,
     'Every decision traces to this triangle. What problem? For whom? Why does it matter? What can go wrong? Understand before you build.'),
    ('Iceberg Principle', 'Visible: 10% (APIs) | Hidden: 90% (Domain, Security, Data, Failures)', ACCENT_GREEN,
     'Users see only the API surface. But 90% of architectural thinking happens below — domain models, security, data ownership, failure modes.'),
    ('Infinity Loop', 'BUILD (Design→Build→Test→Deploy) ↔ OPERATE (Monitor→Detect→Respond→Learn) ↔ IMPROVE', ACCENT_RED,
     'The system is never "done." It cycles continuously: operate → learn → improve → build. Each iteration improves both the system and the process.'),
]
for i, (title, subtitle, color, desc) in enumerate(models):
    y = Inches(1.15 + i * 1.55)
    card(slide, Inches(0.3), y, Inches(12.7), Inches(1.3))
    pill(slide, Inches(0.45), y + Inches(0.1), Inches(2.8), Inches(0.25), color, title, 11)
    tb(slide, Inches(3.4), y + Inches(0.1), Inches(9.0), Inches(0.22), subtitle, size=10, color=ACCENT_BLUE)
    tb(slide, Inches(0.5), y + Inches(0.45), Inches(12.0), Inches(0.7), desc, size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# S24  KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 24, 'Key Takeaways', 'What to remember from this presentation.')

takeaways = [
    ('Claude Code is Your\nEngineering Team', ACCENT_BLUE,
     'Skills, commands, MCP, sub-agents, hooks, and plugins make it a full development platform — not just a chatbot interface.'),
    ('Think First,\nCode Last', ACCENT_GREEN,
     'Scratchpad → Plan → Tasks → Execute → Verify → Reflect. The 6-phase workflow ensures quality and prevents costly scope creep.'),
    ('Human-in-the-Loop\nis Mandatory', ACCENT_ORANGE,
     'Four approval gates. Explicit "APPROVE" required at each. AI executes with agency — but humans make the decisions.'),
    ('Contracts Define\nEverything', ACCENT_PURPLE,
     'Design API specs, event schemas, and data models BEFORE any implementation. Code conforms to contracts — never the reverse.'),
    ('Vertical Slice\nBefore Horizontal', ACCENT_RED,
     'Prove ONE complete flow end-to-end first. Then expand tier-by-tier. Architecture flaws are exponentially cheaper to catch early.'),
    ('Continuous\nImprovement', ACCENT_YELLOW,
     'The Reflect phase ensures every session makes the next one better. Capture learnings. Update conventions. Evolve the workflow always.'),
]
for i, (title, color, desc) in enumerate(takeaways):
    col = i % 3; row = i // 3
    x = Inches(0.25 + col * 4.3); y = Inches(1.2 + row * 3.05)
    card(slide, x, y, Inches(4.05), Inches(2.7))
    tb(slide, x + Inches(0.15), y + Inches(0.12), Inches(3.75), Inches(0.85), title, size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    pill(slide, x + Inches(0.15), y + Inches(1.1), Inches(3.75), Inches(0.02), color, '', 6)
    tb(slide, x + Inches(0.15), y + Inches(1.25), Inches(3.75), Inches(1.3), desc, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# S25  THANK YOU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide); add_title_bar(slide)
rect_bg(slide, Inches(0), Inches(0.05), Inches(0.06), Inches(7.45), ACCENT_BLUE)

tb(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(1.0), 'Thank You', size=68, color=WHITE, bold=True)
rect_bg(slide, Inches(0.6), Inches(2.6), Inches(1.3), Inches(0.035), ACCENT_GREEN)
tb(slide, Inches(0.6), Inches(3.0), Inches(10), Inches(0.45), 'The future of software development is agentic.', size=26, color=LIGHT_GRAY)
tb(slide, Inches(0.6), Inches(3.55), Inches(10), Inches(0.4), 'Claude Code is your engineering team. You are the architect.', size=20, color=MEDIUM_GRAY)

card(slide, Inches(0.6), Inches(4.3), Inches(11.8), Inches(1.8))
mtb(slide, Inches(0.8), Inches(4.4), Inches(11.3), Inches(1.5), [
    'Start today:',
    '  • Install:    curl -fsSL https://claude.ai/install.sh | bash',
    '  • Learn:      docs.anthropic.com/en/docs/claude-code/overview',
    '  • Build:      Create your CLAUDE.md, custom skills, and first agentic workflow',
    '  • Connect:    Add MCP servers for GitHub, Jira, databases, and more',
    '  • Evolve:     Review, reflect, and improve — the cycle never ends',
], size=13, color=LIGHT_GRAY, ls=1.22)

tb(slide, Inches(0.6), Inches(6.5), Inches(11), Inches(0.25),
    'code.claude.com  |  docs.anthropic.com  |  github.com/modelcontextprotocol',
    size=12, color=ACCENT_BLUE)


# ── Save ────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Agentic_Development_with_Claude_Code_v2.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
